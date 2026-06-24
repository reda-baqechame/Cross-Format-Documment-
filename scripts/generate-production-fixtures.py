#!/usr/bin/env python3
"""Write binary fixtures for scripts/production-full-test.mjs.

Run from repo root:
  cd backend && uv run python ../scripts/generate-production-fixtures.py
"""

from __future__ import annotations

import io
from pathlib import Path

OUT = Path(__file__).resolve().parent / "production-fixtures"


def write_pdf() -> None:
    import fitz

    pdf = fitz.open()
    pdf.set_metadata({"title": "Production QA Contract", "author": "DocOS QA"})
    page = pdf.new_page(width=595, height=842)
    page.insert_text((72, 72), "SERVICE AGREEMENT — Random Client 847293", fontsize=14)
    page.insert_text((72, 110), "Total fee: USD 45,000. Contact: billing@client.example", fontsize=11)
    page.insert_text((72, 150), "SSN 987-65-4321 for onboarding verification.", fontsize=11)
    page.insert_text((72, 190), "Signature: ___________________", fontsize=11)
    (OUT / "random-contract.pdf").write_bytes(pdf.tobytes())
    pdf.close()


def write_docx() -> None:
    from docx import Document

    doc = Document()
    doc.core_properties.title = "Random QA Proposal"
    doc.core_properties.author = "DocOS QA"
    doc.add_heading("Proposal for Client 847293", level=1)
    doc.add_paragraph("Scope: document platform rollout with portal and audit.")
    doc.add_paragraph("Payment: USD 45,000 net-30.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Milestone"
    table.cell(0, 1).text = "Amount"
    table.cell(1, 0).text = "Delivery"
    table.cell(1, 1).text = "45000"
    buf = io.BytesIO()
    doc.save(buf)
    (OUT / "random-proposal.docx").write_bytes(buf.getvalue())


def write_xlsx() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Invoice"
    ws.append(["Invoice #", "847293"])
    ws.append(["Amount", 12500])
    ws.append(["Client", "Northwind"])
    buf = io.BytesIO()
    wb.save(buf)
    (OUT / "random-invoice.xlsx").write_bytes(buf.getvalue())


def write_pptx() -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Deck QA 847293"
    slide.placeholders[1].text = "Quarterly review slide content"
    buf = io.BytesIO()
    prs.save(buf)
    (OUT / "random-deck.pptx").write_bytes(buf.getvalue())


def write_rtf() -> None:
    data = rb"{\rtf1\ansi Random RTF document 847293\par Fee USD 45000\par}"
    (OUT / "random-note.rtf").write_bytes(data)


def write_png() -> None:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (240, 120), color=(230, 240, 255)).save(buf, format="PNG")
    (OUT / "random-scan.png").write_bytes(buf.getvalue())


def write_html() -> None:
    html = """<!DOCTYPE html>
<html><head><title>Random HTML 847293</title></head>
<body><h1>Invoice</h1><p>Amount: USD 5000</p><p>Client: Northwind</p></body></html>"""
    (OUT / "random-page.html").write_bytes(html.encode("utf-8"))


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    write_pdf()
    write_docx()
    write_xlsx()
    write_pptx()
    write_rtf()
    write_png()
    write_html()
    print(f"[fixtures] wrote 7 files to {OUT}")


if __name__ == "__main__":
    main()
