"""PPTX adapter (python-pptx).

Maps each slide to a page, each shape's text frame to paragraphs and runs (with
bold/italic/size where set), and pictures to ImageNodes carrying alt text when the
deck provides it. Slides flow through the same node graph as every other format, so a
deck is editable and exportable to DOCX/TXT.
"""

from __future__ import annotations

import io
from datetime import UTC, datetime

from pptx import Presentation
from pptx.util import Emu

from docos.model.document import CanonicalDocument, DocumentMeta
from docos.model.ids import new_doc_id, new_node_id
from docos.model.nodes import ImageNode, PageNode, ParagraphNode, RootNode, RunNode
from docos.services.docengine.interface import FormatAdapter
from docos.storage.blob import BlobStore

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class PptxAdapter(FormatAdapter):
    format_id = "pptx"
    supported_mimes = (_PPTX_MIME,)

    def parse(self, data: bytes, *, blob: BlobStore | None = None) -> CanonicalDocument:
        prs = Presentation(io.BytesIO(data))
        now = datetime.now(UTC)
        root = RootNode(id=new_node_id("root"))
        core = prs.core_properties
        meta = DocumentMeta(
            title=core.title or None,
            author=core.author or None,
            source_format="pptx",
            source_mime=_PPTX_MIME,
            created_at=now,
            modified_at=now,
            page_count=len(prs.slides),
            custom={"last_modified_by": core.last_modified_by, "keywords": core.keywords},
        )
        doc = CanonicalDocument(doc_id=new_doc_id(), root_id=root.id, meta=meta)
        doc.add_node(root)

        width = Emu(prs.slide_width or 0).pt
        height = Emu(prs.slide_height or 0).pt
        for idx, slide in enumerate(prs.slides):
            page = PageNode(
                id=new_node_id("page"),
                parent_id=root.id,
                page_number=idx + 1,
                width=width,
                height=height,
                reading_order=idx,
            )
            root.children.append(page.id)
            doc.add_node(page)
            for shape in slide.shapes:
                if shape.shape_type == 13:  # PICTURE
                    self._add_image(doc, page, shape)
                elif shape.has_text_frame:
                    self._add_text_frame(doc, page, shape)

        if meta.title:
            doc.accessibility.has_doc_title = True
        return doc

    def _add_text_frame(self, doc: CanonicalDocument, page: PageNode, shape) -> None:
        for para in shape.text_frame.paragraphs:
            pnode = ParagraphNode(id=new_node_id(), parent_id=page.id)
            for run in para.runs:
                if not run.text:
                    continue
                font = run.font
                rnode = RunNode(
                    id=new_node_id(),
                    parent_id=pnode.id,
                    text=run.text,
                    bold=bool(font.bold),
                    italic=bool(font.italic),
                    underline=bool(font.underline),
                    font=font.name,
                    size=float(font.size.pt) if font.size else None,
                )
                pnode.children.append(rnode.id)
                doc.add_node(rnode)
            if pnode.children:
                page.children.append(pnode.id)
                doc.add_node(pnode)

    def _add_image(self, doc: CanonicalDocument, page: PageNode, shape) -> None:
        alt = (getattr(shape, "name", None) or "image").strip()
        blob_ref = f"pptx/{page.page_number}/{shape.shape_id}"
        # Pull the picture bytes so the upload route can persist them (parse is sync; persistence
        # is async). Guard against malformed pictures so a bad shape never fails the whole parse.
        img_bytes = b""
        mime: str | None = None
        try:
            image = shape.image
            img_bytes = image.blob or b""
            mime = image.content_type or None
        except Exception:  # noqa: BLE001 - a picture without retrievable bytes degrades gracefully
            img_bytes = b""
        node = ImageNode(
            id=new_node_id("img"),
            parent_id=page.id,
            blob_ref=blob_ref,
            mime=mime,
            alt_text=alt or None,
            attrs={"persisted": False},
        )
        if img_bytes:
            doc._pending_assets[blob_ref] = img_bytes
        page.children.append(node.id)
        doc.add_node(node)

    def render_preview(self, doc: CanonicalDocument, page: int) -> bytes:
        raise NotImplementedError("PptxAdapter.render_preview — slides render in the canvas")

    def export(self, doc: CanonicalDocument, *, target_mime: str) -> bytes:
        from docos.services.docengine.writers.pptx_writer import model_to_pptx

        return model_to_pptx(doc)
