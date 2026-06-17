"""Canonical-model → XLSX / PPTX writers.

Completes the "convert anything to any Office format" matrix: any opened document —
PDF, DOCX, a scan, a spreadsheet — downloads as a real ``.xlsx`` or ``.pptx`` because
every format already lives in one node graph. Redaction is honored through
``run_text`` exactly like the DOCX/PDF/markup writers, so redacted content never
reaches the output.
"""

from __future__ import annotations

import io

from docos.model.document import CanonicalDocument
from docos.model.nodes import AnyNode
from docos.services.docengine.writers.redaction import run_text


def _runs(doc: CanonicalDocument, block: AnyNode) -> list[AnyNode]:
    return [c for c in doc.children_of(block.id) if c.type == "run"]


def _plain(doc: CanonicalDocument, block: AnyNode) -> str:
    return "".join(run_text(doc, r) for r in _runs(doc, block)).strip()


def _table_rows(doc: CanonicalDocument, tnode: AnyNode) -> list[list[str]]:
    rows: list[list[str]] = []
    for row in (n for n in doc.children_of(tnode.id) if n.type == "table_row"):
        cells = [c for c in doc.children_of(row.id) if c.type == "table_cell"]
        rows.append([_plain(doc, cell) for cell in cells])
    return rows


def _sheet_title(name: str, used: set[str]) -> str:
    # Excel sheet names: <=31 chars, no []:*?/\ and must be unique.
    safe = "".join("_" if ch in r"[]:*?/\\" else ch for ch in name)[:31] or "Sheet"
    candidate, i = safe, 1
    while candidate in used:
        suffix = f" {i}"
        candidate = safe[: 31 - len(suffix)] + suffix
        i += 1
    used.add(candidate)
    return candidate


def model_to_xlsx(doc: CanonicalDocument) -> bytes:
    """Each table becomes a worksheet; remaining text goes to a 'Text' sheet."""
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)  # start empty; we add sheets deliberately
    used: set[str] = set()

    for node in doc.walk():
        if node.type == "table":
            rows = _table_rows(doc, node)
            if not rows:
                continue
            ws = wb.create_sheet(title=_sheet_title("Table", used))
            for row in rows:
                ws.append(row)

    text_rows = [
        _plain(doc, n) for n in doc.walk() if n.type in ("heading", "paragraph") and _plain(doc, n)
    ]
    if text_rows:
        ws = wb.create_sheet(title=_sheet_title("Text", used))
        for line in text_rows:
            ws.append([line])

    if not wb.sheetnames:  # empty document — keep the file valid
        wb.create_sheet(title="Sheet1")

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def model_to_pptx(doc: CanonicalDocument) -> bytes:
    """One slide per page (paged formats) or per heading-section (flow formats)."""
    from pptx import Presentation
    from pptx.util import Emu, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]
    margin = Emu(457200)  # 0.5"
    box_w = prs.slide_width - margin * 2
    box_h = prs.slide_height - margin * 2

    sections = _sections(doc)
    if not sections:
        sections = [("", [])]

    for title, lines in sections:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(margin, margin, box_w, box_h)
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        if title:
            p = tf.paragraphs[0]
            p.text = title
            p.font.bold = True
            p.font.size = Pt(24)
            first = False
        for line in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            first = False
        if first:  # nothing added — leave an empty placeholder line
            tf.paragraphs[0].text = ""

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _sections(doc: CanonicalDocument) -> list[tuple[str, list[str]]]:
    """Group top-level content into (title, body-lines) sections for slides.

    A new section starts at each page or heading; tables flatten to tab-joined rows.
    """
    sections: list[tuple[str, list[str]]] = []
    title, lines = "", []

    def flush() -> None:
        nonlocal title, lines
        if title or lines:
            sections.append((title, lines))
        title, lines = "", []

    def walk(node: AnyNode) -> None:
        nonlocal title, lines
        for child in doc.children_of(node.id):
            kind = child.type
            if kind == "page":
                flush()
                walk(child)
                flush()
            elif kind == "heading":
                flush()
                title = _plain(doc, child)
            elif kind == "paragraph":
                text = _plain(doc, child)
                if text:
                    lines.append(text)
            elif kind == "list":
                for item in (c for c in doc.children_of(child.id) if c.type == "list_item"):
                    text = _plain(doc, item)
                    if text:
                        lines.append(f"• {text}")
            elif kind == "table":
                for row in _table_rows(doc, child):
                    lines.append("\t".join(row))
            elif kind == "field":
                name = getattr(child, "field_name", "field")
                lines.append(f"{name}: {getattr(child, 'value', '') or ''}")

    walk(doc.nodes[doc.root_id])
    flush()
    return sections
