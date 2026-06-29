"""Universal canonical-model → PPTX writer.

Rebuilds a real ``.pptx`` from the node graph. If the document already has page
nodes (a PPTX or PDF origin), each page becomes one slide. Otherwise the document
is sectioned into slides by heading — each heading starts a new slide whose body
gathers the following paragraphs/lists — so a TXT or DOCX can be downloaded as a
deck. Tolerant by design: unknown nodes are skipped rather than raising.
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Inches, Pt

from docos.model.document import CanonicalDocument
from docos.model.nodes import AnyNode
from docos.services.docengine.writers.redaction import is_redacted, node_text, run_text

_BLANK_LAYOUT = 6  # the standard blank layout in the default template


def _block_text(doc: CanonicalDocument, block: AnyNode) -> str:
    parts: list[str] = []
    for child in doc.children_of(block.id):
        if child.type == "run":
            parts.append(run_text(doc, child))
        elif child.type == "footnote_reference" and not is_redacted(doc, child.id):
            parts.append(f"[{getattr(child, 'marker', '')}]")
        elif child.type == "unsupported":
            parts.append(f"[unsupported: {getattr(child, 'original_type', 'unknown')}]")
    return "".join(parts)


def _footnote_line(doc: CanonicalDocument, footnote: AnyNode) -> str:
    if is_redacted(doc, footnote.id):
        return ""
    parts: list[str] = []
    direct = _block_text(doc, footnote).strip()
    if direct:
        parts.append(direct)
    for child in doc.children_of(footnote.id):
        if child.type in ("paragraph", "heading", "table_cell"):
            text = _block_text(doc, child).strip()
            if text:
                parts.append(text)
    text = " ".join(parts).strip()
    return f"{getattr(footnote, 'marker', '')}. {text}" if text else ""


def _table_lines(doc: CanonicalDocument, tnode: AnyNode) -> list[str]:
    lines: list[str] = []
    for row in doc.children_of(tnode.id):
        if row.type != "table_row":
            continue
        cells = [_block_text(doc, c) for c in doc.children_of(row.id) if c.type == "table_cell"]
        lines.append(" | ".join(cells))
    return lines


def _add_slide(
    prs: Presentation, title: str, body_lines: list[str], image_bytes: list[bytes] | None = None
) -> None:
    image_bytes = image_bytes or []
    slide = prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])
    # Narrow the text column when there are images so the pictures get the right half of the slide.
    text_width = Inches(5.5) if image_bytes else Inches(9)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), text_width, Inches(6.5))
    tf = box.text_frame
    tf.word_wrap = True
    first = tf.paragraphs[0]
    first.text = title or ""
    first.font.size = Pt(28)
    first.font.bold = True
    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)

    # Stack embedded pictures down the right column (best-effort; skip any that won't load).
    top = Inches(0.5)
    for data in image_bytes:
        try:
            slide.shapes.add_picture(io.BytesIO(data), Inches(6.2), top, height=Inches(1.8))
            top += Inches(2.0)
        except Exception:  # noqa: BLE001 - unreadable/unsupported image is simply omitted
            continue


def model_to_pptx(doc: CanonicalDocument, images: dict[str, bytes] | None = None) -> bytes:
    """Serialise ``doc`` to PPTX. ``images`` maps ``ImageNode.blob_ref`` → bytes; available image
    bytes are embedded as picture shapes, and images without bytes fall back to a text note."""
    images = images or {}
    prs = Presentation()
    top = doc.children_of(doc.root_id)
    pages = [n for n in top if n.type == "page"]

    if pages:
        for page in pages:
            title = ""
            body: list[str] = []
            slide_images: list[bytes] = []
            for child in doc.children_of(page.id):
                if child.type == "heading" and not title:
                    title = _block_text(doc, child)
                elif child.type in ("paragraph", "heading"):
                    text = _block_text(doc, child)
                    if text:
                        body.append(text)
                elif child.type == "list":
                    for item in doc.children_of(child.id):
                        if item.type == "list_item":
                            body.append(f"• {_block_text(doc, item)}")
                elif child.type == "table":
                    body.extend(_table_lines(doc, child))
                elif child.type == "footnote":
                    line = _footnote_line(doc, child)
                    if line:
                        body.append(f"Footnote {line}")
                elif child.type == "unsupported":
                    body.append(f"[unsupported node: {getattr(child, 'original_type', 'unknown')}]")
                elif child.type == "image":
                    if is_redacted(doc, child.id):
                        continue
                    data = images.get(getattr(child, "blob_ref", "") or "")
                    if data:
                        slide_images.append(data)
                    else:
                        body.append(f"[image: {node_text(doc, child) or 'image'}]")
            _add_slide(prs, title or f"Slide {page.page_number}", body, slide_images)
    else:
        # Section a flat document into slides by heading.
        title = doc.meta.title or "Slide 1"
        body = []
        opened = False
        for node in top:
            if node.type == "heading":
                if opened or body:
                    _add_slide(prs, title, body)
                title = _block_text(doc, node) or "Slide"
                body = []
                opened = True
            elif node.type in ("paragraph", "list_item"):
                text = _block_text(doc, node)
                if text:
                    body.append(text)
            elif node.type == "list":
                for item in doc.children_of(node.id):
                    if item.type == "list_item":
                        body.append(f"• {_block_text(doc, item)}")
            elif node.type == "table":
                body.extend(_table_lines(doc, node))
            elif node.type == "footnote":
                line = _footnote_line(doc, node)
                if line:
                    body.append(f"Footnote {line}")
            elif node.type == "unsupported":
                body.append(f"[unsupported node: {getattr(node, 'original_type', 'unknown')}]")
        _add_slide(prs, title, body)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
