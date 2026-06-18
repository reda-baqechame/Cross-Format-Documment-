"""Shared fixtures: sample documents built on the fly + a TestClient over SQLite."""

from __future__ import annotations

import io

import pytest
from fastapi.testclient import TestClient
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker

from docos.db.base import Base
from docos.deps import db_session
from docos.main import create_app


@pytest.fixture
def _sessionmaker(tmp_path, monkeypatch):
    """One in-process SQLite db (no Postgres needed), shared by every client in a test."""
    monkeypatch.setenv("LOCAL_BLOB_DIR", str(tmp_path / "blobs"))
    monkeypatch.setenv(
        "ALLOWED_MIME_TYPES",
        "text/plain,application/pdf,"
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document,"
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet,"
        "application/vnd.openxmlformats-officedocument.presentationml.presentation,"
        "application/rtf,image/png,image/jpeg,image/tiff",
    )

    from docos.api import ratelimit

    ratelimit.reset()

    engine = create_engine(f"sqlite:///{tmp_path / 'test.db'}", future=True)
    Base.metadata.create_all(engine)
    return sessionmaker(bind=engine, expire_on_commit=False)


@pytest.fixture
def make_client(_sessionmaker):
    """Factory for TestClients sharing one db.

    Each call returns a *fresh* client with its own cookie jar — i.e. a distinct anonymous
    session — so tests can prove that two sessions can't see each other's documents.
    """

    def _session():
        s = _sessionmaker()
        try:
            yield s
        finally:
            s.close()

    def _make() -> TestClient:
        app = create_app()
        app.dependency_overrides[db_session] = _session
        return TestClient(app)

    return _make


@pytest.fixture
def client(make_client):
    """A single TestClient (one anonymous session)."""
    return make_client()


@pytest.fixture
def db(_sessionmaker):
    """A raw session for asserting persisted state (jobs, tombstones, ownership)."""
    s = _sessionmaker()
    try:
        yield s
    finally:
        s.close()


@pytest.fixture
def sample_txt_bytes() -> bytes:
    return b"Title line\n\nSecond paragraph with some text.\n\nThird paragraph."


@pytest.fixture
def sample_pdf_bytes() -> bytes:
    import fitz

    pdf = fitz.open()
    pdf.set_metadata({"title": "PDF Test", "author": "Tester"})
    page = pdf.new_page(width=595, height=842)  # A4
    page.insert_text((72, 72), "Hello PDF world", fontsize=14)
    page.insert_text((72, 120), "Second line of text", fontsize=11)
    data = pdf.tobytes()
    pdf.close()
    return data


@pytest.fixture
def sample_docx_bytes() -> bytes:
    from docx import Document

    doc = Document()
    doc.core_properties.title = "Test Doc"
    doc.core_properties.author = "Tester"
    doc.add_heading("A Heading", level=1)
    doc.add_paragraph("A normal paragraph with text.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "r0c0"
    table.cell(1, 1).text = "r1c1"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Sales"
    ws.append(["Region", "Total"])
    ws.append(["North", 1200])
    ws.append(["South", 950])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_bytes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "A bullet of slide content"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_rtf_bytes() -> bytes:
    return rb"{\rtf1\ansi First RTF line\par Second RTF line\par}"


@pytest.fixture
def sample_image_bytes() -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (120, 60), color=(240, 240, 240)).save(buf, format="PNG")
    return buf.getvalue()
