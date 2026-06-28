"""Any-format → Markdown / HTML / CSV export over the canonical model."""

from __future__ import annotations

import io


def test_export_docx_as_markdown_and_html(client, sample_docx_bytes):
    docx_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    doc_id = client.post(
        "/documents", files={"file": ("d.docx", io.BytesIO(sample_docx_bytes), docx_mime)}
    ).json()["doc_id"]

    md = client.get(f"/documents/{doc_id}/export", params={"format": "md"})
    assert md.status_code == 200
    assert md.headers["content-type"].startswith("text/markdown")
    body = md.content.decode()
    assert "# A Heading" in body  # heading became an ATX heading
    assert "A normal paragraph with text." in body

    html = client.get(f"/documents/{doc_id}/export", params={"format": "html"})
    assert html.status_code == 200
    assert "<h1>A Heading</h1>" in html.content.decode()


def test_export_txt_doc_as_csv_rows(client):
    doc_id = client.post(
        "/documents", files={"file": ("d.txt", b"Line one\n\nLine two", "text/plain")}
    ).json()["doc_id"]
    csv_out = client.get(f"/documents/{doc_id}/export", params={"format": "csv"})
    assert csv_out.status_code == 200
    assert csv_out.headers["content-type"].startswith("text/csv")
    rows = csv_out.content.decode().splitlines()
    assert "Line one" in rows[0]


def test_redaction_honored_in_markdown(client):
    doc_id = client.post(
        "/documents", files={"file": ("d.txt", b"keep this\n\nsecret line", "text/plain")}
    ).json()["doc_id"]
    model = client.get(f"/documents/{doc_id}/model").json()["document"]
    secret = next(n["id"] for n in model["nodes"].values() if n.get("text") == "secret line")
    client.post(
        f"/documents/{doc_id}/patches", json={"ops": [{"op": "redact", "target_id": secret}]}
    )

    md = client.get(f"/documents/{doc_id}/export", params={"format": "md"}).content.decode()
    assert "keep this" in md
    assert "secret line" not in md


def test_export_docx_as_xlsx(client, sample_docx_bytes):
    docx_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    xlsx_mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    doc_id = client.post(
        "/documents", files={"file": ("d.docx", io.BytesIO(sample_docx_bytes), docx_mime)}
    ).json()["doc_id"]

    out = client.get(f"/documents/{doc_id}/export", params={"format": "xlsx"})
    assert out.status_code == 200
    assert out.headers["content-type"].startswith(xlsx_mime)
    assert out.content[:2] == b"PK"  # xlsx is a zip

    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(out.content))
    text = " ".join(
        str(c.value) for ws in wb.worksheets for row in ws.iter_rows() for c in row if c.value
    )
    assert "A Heading" in text and "r0c0" in text  # paragraph text + table cell present


def test_export_docx_as_pptx(client, sample_docx_bytes):
    docx_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    pptx_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    doc_id = client.post(
        "/documents", files={"file": ("d.docx", io.BytesIO(sample_docx_bytes), docx_mime)}
    ).json()["doc_id"]

    out = client.get(f"/documents/{doc_id}/export", params={"format": "pptx"})
    assert out.status_code == 200
    assert out.headers["content-type"].startswith(pptx_mime)
    assert out.content[:2] == b"PK"

    from pptx import Presentation

    prs = Presentation(io.BytesIO(out.content))
    text = " ".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "A Heading" in text and "normal paragraph" in text


def test_redaction_honored_in_xlsx_and_pptx(client):
    doc_id = client.post(
        "/documents", files={"file": ("d.txt", b"keep this\n\nsecret line", "text/plain")}
    ).json()["doc_id"]
    model = client.get(f"/documents/{doc_id}/model").json()["document"]
    secret = next(n["id"] for n in model["nodes"].values() if n.get("text") == "secret line")
    client.post(
        f"/documents/{doc_id}/patches", json={"ops": [{"op": "redact", "target_id": secret}]}
    )

    from openpyxl import load_workbook
    from pptx import Presentation

    xlsx = client.get(f"/documents/{doc_id}/export", params={"format": "xlsx"}).content
    wb = load_workbook(io.BytesIO(xlsx))
    xtext = " ".join(
        str(c.value) for ws in wb.worksheets for row in ws.iter_rows() for c in row if c.value
    )
    assert "keep this" in xtext and "secret line" not in xtext

    pptx = client.get(f"/documents/{doc_id}/export", params={"format": "pptx"}).content
    prs = Presentation(io.BytesIO(pptx))
    ptext = " ".join(
        s.text_frame.text for sl in prs.slides for s in sl.shapes if s.has_text_frame
    )
    assert "keep this" in ptext and "secret line" not in ptext


def test_unknown_format_rejected(client):
    doc_id = client.post("/documents", files={"file": ("d.txt", b"hi", "text/plain")}).json()[
        "doc_id"
    ]
    assert client.get(f"/documents/{doc_id}/export", params={"format": "xyz"}).status_code == 400


def test_export_spreadsheet_as_xlsx(client, sample_xlsx_bytes):
    xlsx_mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    doc_id = client.post(
        "/documents", files={"file": ("d.xlsx", io.BytesIO(sample_xlsx_bytes), xlsx_mime)}
    ).json()["doc_id"]
    out = client.get(f"/documents/{doc_id}/export", params={"format": "xlsx"})
    assert out.status_code == 200
    assert out.headers["content-type"].startswith(xlsx_mime)
    assert out.content[:2] == b"PK"  # a real OOXML (zip) container


def test_export_any_doc_as_pptx(client, sample_docx_bytes):
    docx_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    doc_id = client.post(
        "/documents", files={"file": ("d.docx", io.BytesIO(sample_docx_bytes), docx_mime)}
    ).json()["doc_id"]
    out = client.get(f"/documents/{doc_id}/export", params={"format": "pptx"})
    assert out.status_code == 200
    assert out.content[:2] == b"PK"


def test_export_any_doc_as_png(client):
    doc_id = client.post(
        "/documents", files={"file": ("d.txt", b"Render me to an image", "text/plain")}
    ).json()["doc_id"]
    out = client.get(f"/documents/{doc_id}/export", params={"format": "png"})
    assert out.status_code == 200
    assert out.headers["content-type"].startswith("image/png")
    assert out.content[:8] == b"\x89PNG\r\n\x1a\n"
