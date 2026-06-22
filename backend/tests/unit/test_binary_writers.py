"""Universal canonical-model → XLSX / PPTX / PNG writers."""

from __future__ import annotations

import io
from datetime import UTC, datetime

from docx import Document as DocxDocument
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation

from docos.model.document import CanonicalDocument, DocumentMeta
from docos.model.ids import new_doc_id, new_node_id
from docos.model.nodes import ImageNode, PageNode, RootNode
from docos.services.docengine.adapters.docx import DocxAdapter
from docos.services.docengine.adapters.txt import TxtAdapter
from docos.services.docengine.adapters.xlsx import XlsxAdapter
from docos.services.docengine.writers.docx_writer import model_to_docx
from docos.services.docengine.writers.image_writer import model_to_png
from docos.services.docengine.writers.pptx_writer import model_to_pptx
from docos.services.docengine.writers.xlsx_writer import model_to_xlsx


def _tiny_png() -> bytes:
    img = Image.new("RGB", (8, 8), (200, 30, 30))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


def _doc_with_image(blob_ref: str, *, persisted: bool) -> CanonicalDocument:
    """A one-image document under a page node (the shape PDF/PPTX produce)."""
    root = RootNode(id=new_node_id("root"))
    now = datetime.now(UTC)
    doc = CanonicalDocument(
        doc_id=new_doc_id(),
        root_id=root.id,
        meta=DocumentMeta(
            source_format="pdf",
            source_mime="application/pdf",
            created_at=now,
            modified_at=now,
        ),
    )
    doc.add_node(root)
    page = PageNode(
        id=new_node_id("page"), parent_id=root.id, page_number=1, width=300.0, height=300.0
    )
    root.children.append(page.id)
    doc.add_node(page)
    img = ImageNode(
        id=new_node_id("img"),
        parent_id=page.id,
        blob_ref=blob_ref,
        alt_text="a picture",
        attrs={"persisted": persisted},
    )
    page.children.append(img.id)
    doc.add_node(img)
    return doc


def test_model_to_xlsx_rebuilds_table_from_spreadsheet(sample_xlsx_bytes):
    doc = XlsxAdapter().parse(sample_xlsx_bytes)
    wb = load_workbook(io.BytesIO(model_to_xlsx(doc)))
    values = {
        str(cell.value)
        for ws in wb.worksheets
        for row in ws.iter_rows()
        for cell in row
        if cell.value is not None
    }
    assert {"Region", "Total", "North", "1200"}.issubset(values)


def test_model_to_xlsx_puts_loose_text_on_a_sheet():
    doc = TxtAdapter().parse(b"First line\n\nSecond line")
    wb = load_workbook(io.BytesIO(model_to_xlsx(doc)))
    text = {str(c.value) for ws in wb.worksheets for row in ws.iter_rows() for c in row}
    assert "First line" in text and "Second line" in text


def test_model_to_xlsx_always_has_a_sheet():
    doc = TxtAdapter().parse(b"")
    wb = load_workbook(io.BytesIO(model_to_xlsx(doc)))
    assert len(wb.worksheets) >= 1


def test_model_to_pptx_one_slide_per_page(sample_pptx_bytes):
    from docos.services.docengine.adapters.pptx import PptxAdapter

    doc = PptxAdapter().parse(sample_pptx_bytes)
    prs = Presentation(io.BytesIO(model_to_pptx(doc)))
    assert len(prs.slides) >= 1
    all_text = " ".join(
        run.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for run in para.runs
    )
    assert "Slide Title" in all_text


def test_model_to_pptx_sections_flat_doc_by_heading(sample_docx_bytes):
    doc = DocxAdapter().parse(sample_docx_bytes)
    prs = Presentation(io.BytesIO(model_to_pptx(doc)))
    text = " ".join(
        run.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
        for para in shape.text_frame.paragraphs
        for run in para.runs
    )
    assert "A Heading" in text and "normal paragraph" in text


def test_model_to_xlsx_includes_paged_content(sample_pptx_bytes):
    """PPTX/PDF content lives under page nodes; the writer must descend into them."""
    from docos.services.docengine.adapters.pptx import PptxAdapter

    doc = PptxAdapter().parse(sample_pptx_bytes)
    wb = load_workbook(io.BytesIO(model_to_xlsx(doc)))
    text = " ".join(
        str(c.value) for ws in wb.worksheets for row in ws.iter_rows() for c in row if c.value
    )
    assert "Slide Title" in text


def test_model_to_docx_embeds_image_bytes_when_available():
    doc = _doc_with_image("images/abc", persisted=True)
    data = model_to_docx(doc, {"images/abc": _tiny_png()})
    rendered = DocxDocument(io.BytesIO(data))
    assert len(rendered.inline_shapes) == 1  # real picture embedded, not a placeholder


def test_model_to_docx_falls_back_to_placeholder_without_bytes():
    doc = _doc_with_image("images/abc", persisted=True)
    data = model_to_docx(doc)  # no images map
    rendered = DocxDocument(io.BytesIO(data))
    assert len(rendered.inline_shapes) == 0
    body_text = " ".join(p.text for p in rendered.paragraphs)
    assert "[image:" in body_text


def test_model_to_pptx_embeds_picture_shape_when_available():
    doc = _doc_with_image("images/abc", persisted=True)
    prs = Presentation(io.BytesIO(model_to_pptx(doc, {"images/abc": _tiny_png()})))
    pictures = [s for slide in prs.slides for s in slide.shapes if s.shape_type == 13]
    assert len(pictures) == 1


def test_model_to_png_is_a_valid_image():
    doc = TxtAdapter().parse(b"Hello image world\n\nA second paragraph here.")
    png = model_to_png(doc)
    img = Image.open(io.BytesIO(png))
    assert img.format == "PNG"
    assert img.width > 0 and img.height > 0
